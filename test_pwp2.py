from pptx import Presentation

# パワーポイントプレゼンテーションを作成
prs = Presentation()

# スライドを追加
slide_layout = prs.slide_layouts[1]  # スライドのレイアウトを選択（0以外の数値）
slide = prs.slides.add_slide(slide_layout)

# タイトルとサブタイトルを変更
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "新しいタイトル"
subtitle.text = "新しいサブタイトル"

# パワーポイントファイルを保存
prs.save("test2.pptx")

# 作成したパワーポイントファイル内の文字数をカウント
slide_count = len(prs.slides)
text_count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                text_count += len(paragraph.text)

# 文字数を表示
print("スライド数:", slide_count)
print("文字数:", text_count)